#!/usr/bin/env python3
"""
PPTX Hymn Book Parser — Extracts hymns from a PowerPoint hymn book.

Parses a .pptx file where:
  - Each hymn starts with a slide containing a text box with just the hymn number
  - Subsequent slides contain verses in separate text boxes
  - Within a text box, blank paragraphs separate sub-sections (become [SLIDE] breaks)
  - Some hymns have the number embedded on the first line of a text box

Outputs .txt files compatible with ew_tool.py import.

Requires: python-pptx (pip install python-pptx)
"""

import os
import re
import sys
import argparse

try:
    from pptx import Presentation
except ImportError:
    sys.exit("python-pptx is required: pip install python-pptx")

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output", "pptx_book")

META_FENCE = "---"
SLIDE_MARKER = "[SLIDE]"

# Regex: standalone hymn number (1-3 digits, nothing else)
RE_HYMN_NUMBER = re.compile(r'^\d{1,3}$')

# Regex: verse number within packed content (e.g. "01", "02" on its own line)
RE_VERSE_NUMBER = re.compile(r'^\d{1,2}$')


# ---------------------------------------------------------------------------
# Parsing
# ---------------------------------------------------------------------------

def extract_shapes_text(slide):
    """
    Extract text shapes from a slide, categorised as number-only or content.
    Returns (standalone_numbers, content_shapes).
    Each content_shapes entry is a list of (text, is_first_para_number) tuples.
    """
    standalone_numbers = []
    content_shapes = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        full_text = shape.text_frame.text.strip()
        if not full_text:
            continue

        if RE_HYMN_NUMBER.match(full_text):
            standalone_numbers.append(int(full_text))
        else:
            # Collect paragraph-level data for this shape
            paras = []
            for para in shape.text_frame.paragraphs:
                paras.append(para.text)
            content_shapes.append(paras)

    return standalone_numbers, content_shapes


def paras_to_verses(paras, strip_verse_nums=False):
    """
    Convert a list of paragraph strings into verse blocks.
    Blank paragraphs act as separators between verse blocks.
    If strip_verse_nums is True, paragraphs that are just 1-2 digit numbers
    are treated as verse-number labels and stripped.
    Returns a list of verse strings.
    """
    verses = []
    current_lines = []

    for para_text in paras:
        stripped = para_text.strip()

        if not stripped:
            # Blank paragraph = verse separator
            if current_lines:
                verses.append('\n'.join(current_lines))
                current_lines = []
            continue

        if strip_verse_nums and RE_VERSE_NUMBER.match(stripped):
            # This is a verse number label — start a new verse
            if current_lines:
                verses.append('\n'.join(current_lines))
                current_lines = []
            continue

        # Clean up excessive whitespace/tabs in the line
        cleaned = re.sub(r'[\t ]+', ' ', stripped).strip()
        if cleaned:
            current_lines.append(cleaned)

    if current_lines:
        verses.append('\n'.join(current_lines))

    return verses


def parse_pptx(filepath):
    """
    Parse a .pptx hymn book file.
    Returns a list of dicts: {number, verses: [str, ...]}
    Each verse string is the text for one EW slide.
    """
    prs = Presentation(filepath)
    slides = list(prs.slides)

    hymns = []
    current_hymn = None

    for slide_idx, slide in enumerate(slides):
        standalone_nums, content_shapes = extract_shapes_text(slide)

        # Determine if this slide starts a new hymn
        # Priority 1: standalone number textbox
        if standalone_nums:
            for num in standalone_nums:
                # Each standalone number starts a new hymn
                if current_hymn is not None and current_hymn['verses']:
                    hymns.append(current_hymn)
                current_hymn = {
                    'number': num,
                    'number_str': str(num).zfill(2) if num < 100 else str(num),
                    'verses': [],
                    'start_slide': slide_idx,
                }

            # Process content shapes as verses of the current hymn
            if current_hymn is not None:
                for paras in content_shapes:
                    # Check if first paragraph is a small number (verse label)
                    first_stripped = paras[0].strip() if paras else ''
                    has_verse_nums = RE_VERSE_NUMBER.match(first_stripped) and len(paras) > 1

                    verses = paras_to_verses(paras, strip_verse_nums=has_verse_nums)
                    current_hymn['verses'].extend(verses)

        elif content_shapes:
            # No standalone number on this slide
            # Check if any content shape starts with an embedded hymn number
            first_shape_paras = content_shapes[0]
            first_para_text = first_shape_paras[0].strip() if first_shape_paras else ''

            if RE_HYMN_NUMBER.match(first_para_text) and len(first_shape_paras) > 1:
                # Embedded hymn number — first paragraph is the number
                if current_hymn is not None and current_hymn['verses']:
                    hymns.append(current_hymn)

                num = int(first_para_text)
                current_hymn = {
                    'number': num,
                    'number_str': str(num).zfill(2) if num < 100 else str(num),
                    'verses': [],
                    'start_slide': slide_idx,
                }

                # Rest of this shape (after the number) is verse content
                remaining_paras = first_shape_paras[1:]
                has_verse_nums = any(
                    RE_VERSE_NUMBER.match(p.strip()) for p in remaining_paras
                    if p.strip()
                )
                verses = paras_to_verses(remaining_paras, strip_verse_nums=has_verse_nums)
                current_hymn['verses'].extend(verses)

                # Process remaining shapes normally
                for paras in content_shapes[1:]:
                    first_stripped = paras[0].strip() if paras else ''
                    has_vn = RE_VERSE_NUMBER.match(first_stripped) and len(paras) > 1
                    verses = paras_to_verses(paras, strip_verse_nums=has_vn)
                    current_hymn['verses'].extend(verses)
            else:
                # Continuation slide — add verses to current hymn
                if current_hymn is not None:
                    for paras in content_shapes:
                        first_stripped = paras[0].strip() if paras else ''
                        has_vn = RE_VERSE_NUMBER.match(first_stripped) and len(paras) > 1
                        verses = paras_to_verses(paras, strip_verse_nums=has_vn)
                        current_hymn['verses'].extend(verses)

    # Don't forget the last hymn
    if current_hymn is not None and current_hymn['verses']:
        hymns.append(current_hymn)

    return hymns


# ---------------------------------------------------------------------------
# Output
# ---------------------------------------------------------------------------

def hymn_to_lyrics(hymn):
    """Convert a hymn's verses list to lyrics text with [SLIDE] markers."""
    return ('\n' + SLIDE_MARKER + '\n').join(hymn['verses'])


def write_hymn_txt(filepath, hymn_number, lyrics):
    """Write a single hymn to a .txt file in ew_tool.py format."""
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(f'{META_FENCE}\n')
        f.write(f'title: {hymn_number}\n')
        f.write(f'title_sinhala: \n')
        f.write(f'author: \n')
        f.write(f'copyright: \n')
        f.write(f'ccli: \n')
        f.write(f'book_ref: {hymn_number}\n')
        f.write(f'source: pptx_book\n')
        f.write(f'{META_FENCE}\n\n')
        f.write(lyrics)
        f.write('\n')


def export_all(pptx_path, output_dir=None):
    """Parse a PPTX hymn book and export all hymns as .txt files."""
    if output_dir is None:
        output_dir = OUTPUT_DIR

    os.makedirs(output_dir, exist_ok=True)

    print(f"Parsing: {pptx_path}")
    hymns = parse_pptx(pptx_path)
    print(f"Found {len(hymns)} hymns")

    written = 0
    for hymn in hymns:
        number_str = hymn['number_str']
        lyrics = hymn_to_lyrics(hymn)

        if not lyrics.strip():
            continue

        filename = f"hymn_{number_str}.txt"
        filepath = os.path.join(output_dir, filename)

        # Handle duplicate filenames (multiple sections may share numbers)
        counter = 1
        base_path = filepath
        while os.path.exists(filepath):
            name, ext = os.path.splitext(base_path)
            filepath = f"{name}_{counter}{ext}"
            counter += 1

        write_hymn_txt(filepath, number_str, lyrics)
        written += 1

    print(f"Exported {written} hymns to: {output_dir}")
    return hymns


# ---------------------------------------------------------------------------
# Entry Point
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        description="Extract hymns from a PowerPoint hymn book"
    )
    parser.add_argument('pptx_file', help="Path to the .pptx file")
    parser.add_argument('--output', default=None,
                        help=f"Output directory (default: {OUTPUT_DIR})")
    parser.add_argument('--test', type=int, default=0,
                        help="Only show first N hymns (no file output)")
    parser.add_argument('--stats', action='store_true',
                        help="Show statistics about the hymn book")
    args = parser.parse_args()

    if args.output:
        OUTPUT_DIR = args.output

    if args.stats:
        hymns = parse_pptx(args.pptx_file)
        print(f"Total hymns: {len(hymns)}")
        nums = [h['number'] for h in hymns]
        print(f"Number range: {min(nums)} - {max(nums)}")
        verse_counts = [len(h['verses']) for h in hymns]
        print(f"Verses per hymn: min={min(verse_counts)}, max={max(verse_counts)}, "
              f"avg={sum(verse_counts)/len(verse_counts):.1f}")
        # Number frequency
        from collections import Counter
        vc = Counter(verse_counts)
        print(f"Verse count distribution: {sorted(vc.items())}")

    elif args.test > 0:
        hymns = parse_pptx(args.pptx_file)
        for hymn in hymns[:args.test]:
            lyrics = hymn_to_lyrics(hymn)
            print(f"\n{'='*60}")
            print(f"Hymn {hymn['number_str']} ({len(hymn['verses'])} slides)")
            print(f"{'='*60}")
            print(lyrics)

    else:
        export_all(args.pptx_file, args.output)
